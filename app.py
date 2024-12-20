import streamlit as st
from pptx import Presentation
import docx
import PyPDF2
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity

# Load pre-trained model for embeddings
MODEL_NAME = 'all-MiniLM-L6-v2'
model = SentenceTransformer(MODEL_NAME)

# Title and description
st.title("Document Similarity Checker")
st.write("Upload two documents to compare their vector embeddings and compute their similarity.")

# File uploaders
doc1 = st.file_uploader("Upload the first document", type=["txt", "pdf", "docx", "pptx"])
doc2 = st.file_uploader("Upload the second document", type=["txt", "pdf", "docx", "pptx"])

# Defines a function that loads text from a Word file
def word_file_load(file_name):
    doc = docx.Document(file_name)
    fullText = [para.text for para in doc.paragraphs]
    return '\n'.join(fullText)

# Defines a function that loads text from a PowerPoint file
def pptx_file_load(file_name):
    prs = Presentation(file_name)
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                fullText.append(shape.text)
    return '\n'.join(fullText)

def extract_text(file):
    file_content = ""
    file_name = file.name
    if file_name.endswith(".pdf"):
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            file_content += " " + page.extract_text()
    elif file_name.endswith(".txt"):
        file_content = file.read().decode("utf-8")
    elif file_name.endswith(".docx"):
        file_content = word_file_load(file)
    elif file_name.endswith(".pptx"):
        file_content = pptx_file_load(file)
    return file_content

# Process and compute similarity on button click
if st.button("Submit"):
    if doc1 and doc2:
        text1 = extract_text(doc1)
        text2 = extract_text(doc2)

        if text1 and text2:
            # Display the documents
            st.subheader("Document 1 Preview:")
            st.text_area("", text1[:500] + ("..." if len(text1) > 500 else ""), height=200, disabled=True)

            st.subheader("Document 2 Preview:")
            st.text_area("", text2[:500] + ("..." if len(text2) > 500 else ""), height=200, disabled=True)

            # Compute embeddings
            st.write("Computing vector embeddings...")
            embeddings1 = model.encode([text1])[0]
            embeddings2 = model.encode([text2])[0]

            # Compute cosine similarity
            similarity = cosine_similarity([embeddings1], [embeddings2])[0][0]

            # Display similarity score
            st.subheader("Similarity Score")
            st.metric(label="Cosine Similarity", value=f"{similarity:.4f}")
        else:
            st.warning("Please ensure both documents are valid and not empty.")
    else:
        st.warning("Please upload both documents to proceed.")
